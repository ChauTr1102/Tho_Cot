#!/usr/bin/env python3
"""
Gen bản PowerPoint pitch deck cho Team Thợ Cốt — AI Cross-Border Hackathon 2026.
Brand kit bóc từ aiglobal.dev (xem assets/style-guide.md).
Chạy:  python3 build_pptx.py   ->  tho-cot-pitch-template.pptx
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image, ImageDraw, ImageFilter, ImageChops

_HERE = os.path.dirname(os.path.abspath(__file__))
LOGO_EVENT = os.path.join(_HERE, "brand", "logo-event.png")   # wordmark ngang của sự kiện
MASCOT = os.path.join(_HERE, "brand", "hero-mascot.png")      # linh vật nón lá (cắt từ hero banner)

# ---- Brand tokens ----
BG      = RGBColor(0x00, 0x17, 0x08)
FG      = RGBColor(0xF4, 0xFA, 0xF5)
ACCENT  = RGBColor(0x7E, 0xF9, 0x62)
PRIMARY = RGBColor(0x35, 0xEA, 0x52)
GOLD    = RGBColor(0xFC, 0xBB, 0x00)
CTA     = RGBColor(0xFE, 0x6E, 0x00)
MUTEDFG = RGBColor(0xA1, 0xC1, 0xA7)
CARD    = RGBColor(0x00, 0x22, 0x0E)
BORDER  = RGBColor(0x24, 0x4A, 0x2F)
GOLDBRD = RGBColor(0x6B, 0x53, 0x12)

DISP = "Space Grotesk"   # tiêu đề (fallback nếu máy chưa cài)
BODY = "Inter"           # nội dung

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
LX = 0.92                 # lề trái chuẩn
BG_PNG = os.path.join(_HERE, "brand", "bg.png")


def make_bg(path=BG_PNG, W=2560, H=1440):
    """Nền giống HTML: xanh rừng + lưới mờ + glow xanh/gold."""
    img = Image.new("RGB", (W, H), (0x00, 0x17, 0x08))
    glow = Image.new("RGB", (W, H), (0, 0, 0))
    gd = ImageDraw.Draw(glow)

    def blob(cx, cy, r, col):
        gd.ellipse([cx - r, cy - r, cx + r, cy + r], fill=col)

    blob(int(W * 0.04), int(-H * 0.06), int(W * 0.30), (16, 58, 28))   # xanh trên-trái
    blob(int(W * 0.99), int(H * 1.04), int(W * 0.32), (10, 48, 32))    # xanh dưới-phải
    blob(int(W * 0.88), int(H * 0.30), int(W * 0.15), (34, 26, 3))     # gold nhạt
    glow = glow.filter(ImageFilter.GaussianBlur(190))
    img = ImageChops.add(img, glow)
    grid = Image.new("RGBA", (W, H), (0, 0, 0, 0))
    gr = ImageDraw.Draw(grid)
    step = 104
    for x in range(0, W, step):
        gr.line([(x, 0), (x, H)], fill=(126, 249, 98, 9), width=1)
    for y in range(0, H, step):
        gr.line([(0, y), (W, y)], fill=(126, 249, 98, 9), width=1)
    img = Image.alpha_composite(img.convert("RGBA"), grid).convert("RGB")
    img.save(path)


make_bg()


def slide():
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = BG
    if os.path.exists(BG_PNG):
        s.shapes.add_picture(BG_PNG, Inches(0), Inches(0), width=prs.slide_width, height=prs.slide_height)
    return s


def text(s, x, y, w, h, paras, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    """paras: list of dict{runs:[(txt,font,size,bold,color,italic)], after:pt, line:float}"""
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, p in enumerate(paras):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.alignment = align
        para.space_before = Pt(0)
        para.space_after = Pt(p.get("after", 6))
        if "line" in p:
            para.line_spacing = p["line"]
        for (txt, font, size, bold, color, *rest) in p["runs"]:
            r = para.add_run()
            r.text = txt
            r.font.name = font
            r.font.size = Pt(size)
            r.font.bold = bold
            r.font.color.rgb = color
            r.font.italic = rest[0] if rest else False
    return tb


def card(s, x, y, w, h, fill=CARD, line=BORDER, lw=1.0, radius=0.06):
    shp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    try:
        shp.adjustments[0] = radius
    except Exception:
        pass
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = line; shp.line.width = Pt(lw)
    shp.shadow.inherit = False
    return shp


def bar(s, x, y, w, h, color=ACCENT):
    shp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid(); shp.fill.fore_color.rgb = color
    shp.line.fill.background(); shp.shadow.inherit = False
    return shp


def kicker(s, num, label, y=1.55):
    text(s, LX, y - 0.14, 9, 0.4, [{"runs": [
        (f"{num:02d}", DISP, 14, True, ACCENT), ("   ·   ", DISP, 13, False, MUTEDFG),
        (label.upper(), DISP, 11.5, True, MUTEDFG)]}])


def chrome(s, n):
    # brand top-left
    text(s, LX, 0.42, 6, 0.4, [{"runs": [("THỢ CỐT", DISP, 12, True, FG), (".", DISP, 12, True, ACCENT),
         ("   ·   AI Cross-Border Hackathon 2026", BODY, 10, False, MUTEDFG)]}])
    # logo sự kiện góc trên phải (ngữ cảnh: người tham gia)
    if os.path.exists(LOGO_EVENT):
        s.shapes.add_picture(LOGO_EVENT, Inches(11.05), Inches(0.36), height=Inches(0.24))
    # page number bottom-right
    text(s, 11.4, 6.95, 1.5, 0.4,
         [{"runs": [(f"{n:02d}", DISP, 11, True, FG), (f" / 08", BODY, 11, False, MUTEDFG)]}],
         align=PP_ALIGN.RIGHT)


def card_text(shp, paras, anchor=MSO_ANCHOR.TOP, ml=0.22, mt=0.18):
    tf = shp.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = Inches(ml); tf.margin_right = Inches(ml)
    tf.margin_top = Inches(mt); tf.margin_bottom = Inches(mt)
    for i, p in enumerate(paras):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.alignment = p.get("align", PP_ALIGN.LEFT)
        para.space_after = Pt(p.get("after", 4)); para.space_before = Pt(0)
        if "line" in p: para.line_spacing = p["line"]
        for (txt, font, size, bold, color, *rest) in p["runs"]:
            r = para.add_run(); r.text = txt
            r.font.name = font; r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color
            r.font.italic = rest[0] if rest else False


PH = MUTEDFG  # placeholder color

# ============ SLIDE 1 · COVER (editorial competition cover, no motif) ============
s = slide()  # cover không dùng chrome
# nhà tài trợ (2 hàng, căn giữa, trên cùng)
text(s, 0.5, 0.34, 12.33, 0.3, [{"runs": [
    ("TỔ CHỨC BỞI  ", BODY, 8, True, MUTEDFG), ("DNES · Ecomdy        ", DISP, 12, True, FG),
    ("ĐỒNG TỔ CHỨC  ", BODY, 8, True, MUTEDFG), ("Printway · BurgerPrints        ", DISP, 12, True, FG),
    ("POWERED BY  ", BODY, 8, True, MUTEDFG), ("BytePlus", DISP, 12, True, ACCENT)]}], align=PP_ALIGN.CENTER)
text(s, 0.5, 0.66, 12.33, 0.3, [{"runs": [
    ("TÀI TRỢ VÀNG  ", BODY, 8, True, MUTEDFG), ("Kalodata        ", DISP, 12, True, GOLD),
    ("TÀI TRỢ ĐỒNG  ", BODY, 8, True, MUTEDFG), ("Wealify        ", DISP, 12, True, FG),
    ("CROWDFUNDING  ", BODY, 8, True, MUTEDFG),
    ("PG Prints · GKE · Innovark · LianLian · Fristify", DISP, 12, True, FG)]}], align=PP_ALIGN.CENTER)
bar(s, LX, 1.08, 11.5, 0.012, BORDER)
# lead (trái): logo sự kiện + tên sản phẩm lớn + mô tả
if os.path.exists(LOGO_EVENT):
    s.shapes.add_picture(LOGO_EVENT, Inches(LX), Inches(2.15), height=Inches(0.5))
text(s, LX, 2.95, 7.55, 1.7, [{"runs": [("[Tên sản phẩm] ", DISP, 40, True, FG),
     ("— [tagline ngắn]", DISP, 40, True, ACCENT)], "line": 1.02}])
text(s, LX, 4.8, 7.55, 1.2, [{"runs": [
    ("[Mô tả 1–2 câu: sản phẩm là gì, giải quyết vấn đề gì, cho ai, bằng AI như thế nào.]",
     BODY, 14, False, RGBColor(0xc8, 0xdd, 0xca), True)], "line": 1.4}])
# cột thông số (phải) + đường kẻ dọc
bar(s, 9.55, 2.2, 0.014, 4.0, BORDER)
_specs = [("THỜI GIAN", "21–22.08.2026", FG), ("ĐỊA ĐIỂM", "Hà Nội", FG),
          ("HẠNG MỤC", "[chọn nhánh]", FG), ("ĐỘI DỰ THI", "Team Thợ Cốt.", FG),
          ("POWERED BY", "BytePlus", ACCENT)]
_sy = 2.32
for _k, _v, _c in _specs:
    text(s, 9.9, _sy, 3.1, 0.7, [
        {"runs": [(_k, BODY, 9, True, MUTEDFG)], "after": 3},
        {"runs": [(_v, DISP, 15, True, _c)]}])
    _sy += 0.78

# ============ SLIDE 2 · PROBLEM (narrative + stat rows) ============
s = slide(); chrome(s, 2)
kicker(s, 1, "Vấn đề")
text(s, LX, 2.0, 6.7, 1.6, [{"runs": [
    ("[Một câu tuyên bố nỗi đau — VD: “Một lối đi cho tất cả.”]", DISP, 30, True, FG, True)], "line": 1.06}])
text(s, LX, 3.7, 6.3, 0.9, [{"runs": [
    ("[2–3 dòng: ai đang gặp vấn đề gì, vì sao cách làm hiện tại thất bại.]",
     BODY, 14, False, RGBColor(0xc8, 0xdd, 0xca), True)], "line": 1.4}])
text(s, LX, 4.7, 3.5, 1.1, [{"runs": [("[XX%]", DISP, 52, True, ACCENT)]}])
text(s, LX + 3.5, 5.02, 2.9, 1.0, [{"runs": [
    ("[điều gì đang tệ — VD: bỏ cuộc trước khi hoàn thành]", BODY, 14, False, FG, True)], "line": 1.3}])
_rx, _rw, _ry = 8.55, 4.2, 2.35
for _i, (_big, _lbl, _col, _sm) in enumerate([
        ("[<10%]", "[số liệu hỗ trợ 1]", ACCENT, False),
        ("[~50%]", "[số liệu hỗ trợ 2]", GOLD, False),
        ("[Insight]", "[giải pháp lý tưởng đã có nhưng quá đắt để nhân rộng]", FG, True)]):
    if _i > 0:
        bar(s, _rx, _ry, _rw, 0.01, BORDER); _ry += 0.22
    text(s, _rx, _ry, _rw, 0.7, [{"runs": [(_big, DISP, 18 if _sm else 34, True, _col, _sm)]}])
    _ry += 0.5 if _sm else 0.7
    text(s, _rx, _ry, _rw, 0.6, [{"runs": [(_lbl, BODY, 12.5, False, FG if _sm else MUTEDFG, True)], "line": 1.25}])
    _ry += 0.52

# ============ SLIDE 3 · SOLUTION (feature đánh số, không card) ============
s = slide(); chrome(s, 3)
kicker(s, 2, "Giải pháp")
text(s, LX, 2.0, 11.2, 1.3, [{"runs": [
    ("[Sản phẩm] ", DISP, 32, True, ACCENT), ("[làm gì, cho ai, bằng AI]", DISP, 32, True, FG, True)], "line": 1.05}])
for _i, (_idx, _t, _d) in enumerate([
        ("01", "[Tính năng lõi 1]", "[Giá trị cụ thể, đo được nếu có]"),
        ("02", "[Tính năng lõi 2]", "[Giá trị cụ thể]"),
        ("03", "[Tính năng lõi 3]", "[Giá trị cụ thể]")]):
    _x = LX + _i * 3.92
    if _i > 0:
        bar(s, _x - 0.2, 3.95, 0.012, 2.0, BORDER)
    text(s, _x, 3.9, 3.5, 0.7, [{"runs": [(_idx, DISP, 34, True, ACCENT)]}])
    text(s, _x, 4.62, 3.5, 0.5, [{"runs": [(_t, DISP, 18, True, FG, True)]}])
    text(s, _x, 5.12, 3.5, 1.0, [{"runs": [(_d, BODY, 13, False, MUTEDFG, True)], "line": 1.4}])

# ============ SLIDE 4 · DEMO ============
s = slide(); chrome(s, 4)
kicker(s, 3, "Demo")
text(s, LX, 2.0, 10, 0.9, [{"runs": [("Thấy là tin.", DISP, 34, True, FG)]}])
c = card(s, LX, 3.05, 11.5, 3.4, fill=CARD)
card_text(c, [
    {"runs": [("[ Chèn ảnh / video demo sản phẩm ]", DISP, 22, True, FG)], "align": PP_ALIGN.CENTER, "after": 8},
    {"runs": [("Screenshot flow chính · hoặc link demo trực tiếp", BODY, 14, False, MUTEDFG, True)],
     "align": PP_ALIGN.CENTER}], anchor=MSO_ANCHOR.MIDDLE)

# ============ SLIDE 5 · TECH ============
s = slide(); chrome(s, 5)
kicker(s, 4, "Công nghệ")
text(s, LX, 2.0, 11, 0.9, [{"runs": [("Kiến trúc gọn, AI đúng chỗ.", DISP, 32, True, FG)]}])
for _i, (_a, _b, _mid) in enumerate([("[Input]", "[Dữ liệu / yêu cầu từ doanh nghiệp]", False),
                                     ("Lõi AI", "Claude API · [+ BytePlus]", True),
                                     ("[Output]", "[Kết quả cho người dùng]", False)]):
    _x = LX + _i * 4.15
    c = card(s, _x, 3.2, 3.5, 1.5, fill=RGBColor(0x00, 0x1E, 0x0C), line=PRIMARY if _mid else BORDER)
    card_text(c, [
        {"runs": [(_a, DISP, 18, True, ACCENT if _mid else FG)], "align": PP_ALIGN.CENTER, "after": 5},
        {"runs": [(_b, BODY, 12.5, False, MUTEDFG, True)], "align": PP_ALIGN.CENTER, "line": 1.25}],
        anchor=MSO_ANCHOR.MIDDLE)
    if _i < 2:
        text(s, _x + 3.5, 3.5, 0.65, 1, [{"runs": [("→", DISP, 26, True, ACCENT)]}], align=PP_ALIGN.CENTER)
_chips = ["Claude API", "BytePlus", "[Backend]", "[Frontend]", "[Database]", "[Deploy]"]
_cx = LX
for _i, _ch in enumerate(_chips):
    _w = 0.32 + len(_ch) * 0.115
    c = card(s, _cx, 5.15, _w, 0.5, fill=RGBColor(0x01, 0x2A, 0x14), line=PRIMARY if _i == 0 else BORDER, radius=0.5)
    card_text(c, [{"runs": [(_ch, BODY, 12.5, False, ACCENT if _i == 0 else FG)], "align": PP_ALIGN.CENTER}],
              anchor=MSO_ANCHOR.MIDDLE, ml=0.1, mt=0.02)
    _cx += _w + 0.18
text(s, LX, 6.05, 11.5, 0.6, [{"runs": [
    ("💡 Mẹo: ghép BytePlus để “ăn” +5đ thưởng sponsor (xem docs/thong-tin-cuoc-thi.md).",
     BODY, 13, False, MUTEDFG, True)]}])

# ============ SLIDE 6 · IMPACT (hero stat + rows) ============
s = slide(); chrome(s, 6)
kicker(s, 5, "Tác động & Thị trường")
text(s, LX, 2.0, 6.6, 1.6, [{"runs": [("Vì sao ", DISP, 30, True, FG), ("cross-border", DISP, 30, True, ACCENT),
     (", vì sao bây giờ.", DISP, 30, True, FG)], "line": 1.08}])
text(s, LX, 3.85, 3.5, 1.2, [{"runs": [("[$X]", DISP, 52, True, ACCENT)]}])
text(s, LX + 3.5, 4.17, 2.9, 1.0, [{"runs": [("[Quy mô thị trường / TAM có thể chạm tới]", BODY, 14, False, FG, True)], "line": 1.3}])
text(s, LX, 5.55, 6.2, 0.9, [{"runs": [
    ("[1 câu: TMĐT xuyên biên giới VN đang bùng nổ + AI đủ rẻ/mạnh = thời điểm vàng.]",
     BODY, 13.5, False, MUTEDFG, True)], "line": 1.4}])
_rx, _rw, _ry = 8.55, 4.2, 2.35
for _i, (_big, _lbl, _col) in enumerate([("[Nx]", "[nhanh hơn / tiết kiệm so với cách cũ]", ACCENT),
                                         ("[+Y%]", "[tăng doanh thu / hiệu quả cho khách hàng]", GOLD),
                                         ("[Z]", "[chỉ số thứ 3 — người dùng, đơn hàng…]", ACCENT)]):
    if _i > 0:
        bar(s, _rx, _ry, _rw, 0.01, BORDER); _ry += 0.22
    text(s, _rx, _ry, _rw, 0.7, [{"runs": [(_big, DISP, 32, True, _col)]}]); _ry += 0.66
    text(s, _rx, _ry, _rw, 0.5, [{"runs": [(_lbl, BODY, 12.5, False, MUTEDFG, True)], "line": 1.25}]); _ry += 0.54

# ============ SLIDE 7 · ROADMAP (timeline) ============
s = slide(); chrome(s, 7)
kicker(s, 6, "Mô hình & Lộ trình")
text(s, LX, 2.0, 11.2, 0.9, [{"runs": [("Từ nguyên mẫu đến thị trường.", DISP, 32, True, FG)]}])
text(s, LX, 2.95, 11.2, 0.6, [{"runs": [
    ("[Mô hình doanh thu 1 dòng — VD: SaaS theo tháng / phí theo giao dịch]",
     BODY, 14, False, RGBColor(0xc8, 0xdd, 0xca), True)]}])
_ty = 4.2
bar(s, LX + 0.1, _ty, 11.0, 0.02, BORDER)
for _i, (_n, _b, _d, _hl) in enumerate([
        ("BƯỚC 1", "Hackathon", "Nguyên mẫu AI chạy được (48h)", False),
        ("BƯỚC 2", "POC", "Thử nghiệm với 1 doanh nghiệp thật", False),
        ("BƯỚC 3", "DNES Incubation", "Top 10 · mentor · nhà đầu tư", True),
        ("BƯỚC 4", "Ra thị trường", "Khách hàng đầu tiên", False)]):
    _x = LX + _i * 2.95
    _dot = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(_x), Inches(_ty - 0.09), Inches(0.2), Inches(0.2))
    _dot.fill.solid(); _dot.fill.fore_color.rgb = BG
    _dot.line.color.rgb = GOLD if _hl else ACCENT; _dot.line.width = Pt(2.2); _dot.shadow.inherit = False
    text(s, _x, _ty + 0.26, 2.7, 0.3, [{"runs": [(_n, DISP, 11, True, GOLD if _hl else ACCENT)]}])
    text(s, _x, _ty + 0.58, 2.7, 0.4, [{"runs": [(_b, DISP, 17, True, GOLD if _hl else FG)]}])
    text(s, _x, _ty + 1.0, 2.7, 0.9, [{"runs": [(_d, BODY, 12.5, False, MUTEDFG, True)], "line": 1.3}])

# ============ SLIDE 8 · TEAM + THANK YOU ============
s = slide(); chrome(s, 8)
kicker(s, 7, "Đội ngũ")
text(s, LX, 2.0, 11.5, 0.9, [{"runs": [("Team ", DISP, 34, True, FG),
     ("Thợ Cốt", DISP, 34, True, ACCENT), (".", DISP, 34, True, FG)]}])
for _i, (_av, _nm, _role) in enumerate([("T", "[Tên]", "[Vai trò]"), ("H", "[Tên]", "[Vai trò]"),
                                        ("C", "[Tên]", "[Vai trò]"), ("+", "[Tên]", "[Vai trò]")]):
    _x = LX + _i * 2.98
    _sq = card(s, _x, 3.35, 0.85, 0.85, fill=RGBColor(0x02, 0x3A, 0x1C), line=BORDER, radius=0.22)
    card_text(_sq, [{"runs": [(_av, DISP, 24, True, ACCENT)], "align": PP_ALIGN.CENTER}],
              anchor=MSO_ANCHOR.MIDDLE, ml=0, mt=0)
    text(s, _x, 4.4, 2.6, 0.4, [{"runs": [(_nm, DISP, 16, True, FG)]}])
    text(s, _x, 4.78, 2.6, 0.4, [{"runs": [(_role, BODY, 12, False, MUTEDFG, True)]}])
btn = card(s, LX, 5.78, 3.0, 0.68, fill=CTA, line=CTA, radius=0.5)
card_text(btn, [{"runs": [("Cảm ơn & Q/A  →", DISP, 16, True, RGBColor(0x2A, 0x0E, 0x00))],
                 "align": PP_ALIGN.CENTER}], anchor=MSO_ANCHOR.MIDDLE)
text(s, LX + 3.35, 5.9, 8, 0.6, [{"runs": [("[email · SĐT · link demo]", BODY, 15, False, MUTEDFG, True)]}],
     anchor=MSO_ANCHOR.MIDDLE)

out = "tho-cot-pitch-template.pptx"
prs.save(out)
print("Saved:", out, "·", len(prs.slides._sldIdLst), "slides")
